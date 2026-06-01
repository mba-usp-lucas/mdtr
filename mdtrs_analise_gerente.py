"""
Análise MDTR/MDTRS por Gerente — Alcon Commercial Intelligence
================================================================
Script reutilizável: troque o arquivo de entrada e regenere o deck.

USO:
    python mdtrs_analise_gerente.py <caminho_excel.xlsx>

ou edite as variáveis no bloco CONFIG abaixo e rode:
    python mdtrs_analise_gerente.py

VISÃO POR GERENTE (ao invés de 76 slides por setor):
    1 slide por gerente (~7 slides) com:
      - KPIs agregados (YTD, Mai-Proj, Δ vs Abr, Δ S3 vs S2)
      - Top 5 SETORES com MAIOR ALTA (Mai-Proj vs Abr)
        Para cada setor, o PDV de MAIOR impacto positivo
      - Top 5 SETORES OFENSORES (menor performance ou queda)
        Para cada setor, o PDV de MAIOR impacto negativo

ENTRADA esperada: arquivo XLSX com aba "TOTAL" contendo:
- Colunas dimensionais: PROVEDOR, DISTRIBUIDOR AJUSTADO, Bandeira, PDV, CNPJ,
  TIPO_INFORMACAO, TIPO_PDV, ASSOCIACAO, BRICK, POOL, Chave, DESC_BRICK,
  ENDERECO, CIDADE, UF, GERENTE, SETOR_COD, SETOR_NOME, FCC, GRUPO_2,
  GRUPO_MARCA, MARCA, APRESENTACAO
- Colunas mensais: AAAAMM (ex: 202504 = abr/25 ... 202604 = abr/26)
- Mês corrente dividido em semanas + Proj: 202605-S1, 202605-S2, 202605-S3, 202605-Proj
- Coluna YTD-Total

SAÍDA: PPTX com 1 slide capa, 2 slides de produto e 1 slide por gerente.

Autor: Commercial Intelligence — Alcon Brasil
"""

import os
import sys
import io
import warnings
from datetime import datetime

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

warnings.filterwarnings("ignore")

# ============================================================
# CONFIG — edite aqui se quiser rodar sem argumentos
# ============================================================
DEFAULT_INPUT = "/mnt/user-data/uploads/MDTRS_FV_TOTAL.xlsx"
DEFAULT_OUTPUT = "/mnt/user-data/outputs/MDTRS_Analise_Gerente.pptx"
DEFAULT_SHEET = "TOTAL"

# Tipo de informação a analisar (SO = sell-out / SI_NP = sell-in)
TIPO_INFO_ANALISE = "SO"

# Quantos setores mostrar como "destaque" e "ofensor" em cada slide
TOP_SETORES = 5

# Threshold de variação (em unidades) para considerar PDV em destaque
THRESHOLD_VARIACAO = 40

# Excluir o setor "NÃO VISITADO" dos slides individuais?
EXCLUIR_NAO_VISITADO = True

# Mês corrente (semana 1, 2, 3, Proj)
MES_CORRENTE = "202605"  # Maio/2026

# Mês anterior (para comparação Mai vs Abr)
MES_ANTERIOR = "202604"  # Abr/2026

# Período YTD (para análise de produto)
YTD_COLS = ["202601", "202602", "202603", "202604"]  # jan-abr/26
# YTD anterior (mesmo período do ano anterior NÃO disponível; usaremos set-dez/25)
YTD_ANT_COLS = ["202509", "202510", "202511", "202512"]
YTD_LABEL = "YTD 2026 (Jan-Abr)"
YTD_ANT_LABEL = "Set-Dez 2025"

# Paleta de cores — Midnight Executive
NAVY = "1E2761"
ICE_BLUE = "CADCFC"
WHITE = "FFFFFF"
ACCENT_GREEN = "2E7D32"  # alta
ACCENT_RED = "C62828"   # queda
GRAY = "595959"
LIGHT_GRAY = "E8E8E8"
GOLD = "C9A227"

# ============================================================
# UTILS
# ============================================================
def fmt_num(v, casas=0):
    """Formata número no padrão brasileiro."""
    if v is None or (isinstance(v, float) and (np.isnan(v) or np.isinf(v))):
        return "—"
    if casas == 0:
        return f"{int(round(v)):,}".replace(",", ".")
    s = f"{v:,.{casas}f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return s


def fmt_pct(v, casas=1):
    if v is None or (isinstance(v, float) and (np.isnan(v) or np.isinf(v))):
        return "—"
    return f"{v:+.{casas}f}%".replace(".", ",")


def hex_to_rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ============================================================
# DATA LOADER
# ============================================================
def load_data(path, sheet=DEFAULT_SHEET):
    """Lê o XLSX, retornando dataframe com todos os tipos corretos.

    Estratégia para arquivos grandes:
      1) Se houver um pickle de cache mais novo que o XLSX, usa o cache.
         O cache fica em ~/.mdtrs_cache (gravável) — não polui o diretório do XLSX.
      2) Caso contrário, tenta engines em ordem: calamine → openpyxl.
      3) Cache é salvo automaticamente para acelerar reexecuções.
    """
    cache_dir = os.path.join(os.path.expanduser("~"), ".mdtrs_cache")
    os.makedirs(cache_dir, exist_ok=True)
    safe_name = os.path.basename(path).replace(os.sep, "_") + ".pkl"
    cache_path = os.path.join(cache_dir, safe_name)

    if (os.path.exists(cache_path) and
            os.path.getmtime(cache_path) >= os.path.getmtime(path)):
        print(f"[1/5] Lendo cache {cache_path} ...")
        df = pd.read_pickle(cache_path)
    else:
        print(f"[1/5] Lendo {os.path.basename(path)} ...")
        try:
            df = pd.read_excel(path, sheet_name=sheet, engine="calamine")
        except Exception:
            df = pd.read_excel(path, sheet_name=sheet, engine="openpyxl")
        try:
            df.to_pickle(cache_path)
            print(f"      Cache salvo em {cache_path}")
        except Exception as e:
            print(f"      (aviso: não foi possível salvar cache: {e})")

    # Garante numéricos
    period_cols = [c for c in df.columns if str(c).startswith("20") or str(c).startswith("YTD")]
    for c in period_cols:
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    print(f"      Linhas: {len(df):,} | Colunas: {len(df.columns)}")
    return df


# ============================================================
# ANÁLISES
# ============================================================
def analise_gerente(df_gerente, top_n=TOP_SETORES, mes_corrente=MES_CORRENTE, mes_anterior=MES_ANTERIOR):
    """Analisa um gerente: KPIs agregados + ranking de setores (alta/queda)
       + PDV de maior impacto dentro de cada setor."""
    meses_2026 = ["202601", "202602", "202603", "202604"]
    col_proj = f"{mes_corrente}-Proj"
    col_s1 = f"{mes_corrente}-S1"
    col_s2 = f"{mes_corrente}-S2"
    col_s3 = f"{mes_corrente}-S3"

    # KPIs agregados do gerente
    ytd = sum(df_gerente[m].sum() for m in meses_2026)
    proj_mai = df_gerente[col_proj].sum()
    abr = df_gerente[mes_anterior].sum()
    var_mai_abr = proj_mai - abr
    var_mai_abr_pct = (var_mai_abr / abr * 100) if abr > 0 else 0
    s2 = df_gerente[col_s2].sum()
    s3 = df_gerente[col_s3].sum()
    var_s3_s2 = s3 - s2
    var_s3_s2_pct = (var_s3_s2 / s2 * 100) if s2 > 0 else 0

    # Ranking por setor
    setor_grp = df_gerente.groupby("SETOR_NOME").agg(
        abr=(mes_anterior, "sum"),
        mai_proj=(col_proj, "sum"),
        s2=(col_s2, "sum"),
        s3=(col_s3, "sum"),
    )
    setor_grp["var_unid"] = setor_grp["mai_proj"] - setor_grp["abr"]
    setor_grp["var_pct"] = np.where(setor_grp["abr"] > 0,
                                     setor_grp["var_unid"] / setor_grp["abr"] * 100, np.nan)
    setor_grp = setor_grp.reset_index()

    # Para cada setor, identificar o PDV de maior impacto (positivo ou negativo)
    pdv_grp = df_gerente.groupby(["SETOR_NOME", "CNPJ", "PDV", "CIDADE", "UF"]).agg(
        abr=(mes_anterior, "sum"),
        mai_proj=(col_proj, "sum"),
    ).reset_index()
    pdv_grp["var_unid"] = pdv_grp["mai_proj"] - pdv_grp["abr"]

    # Para cada setor, identificar a MARCA de maior impacto (qual produto puxa o setor)
    marca_setor_grp = df_gerente.groupby(["SETOR_NOME", "MARCA"]).agg(
        abr=(mes_anterior, "sum"),
        mai_proj=(col_proj, "sum"),
    ).reset_index()
    marca_setor_grp["var_unid"] = marca_setor_grp["mai_proj"] - marca_setor_grp["abr"]

    def top_pdv(setor, direcao):
        sub = pdv_grp[pdv_grp["SETOR_NOME"] == setor]
        if len(sub) == 0:
            return None
        if direcao == "alta":
            r = sub.loc[sub["var_unid"].idxmax()]
        else:
            r = sub.loc[sub["var_unid"].idxmin()]
        # Sufixo do CNPJ para diferenciar lojas da mesma rede
        nome = str(r["PDV"]).split(" - ")[0][:28]
        try:
            nome = f"{nome} ({str(int(r['CNPJ']))[-4:]})"
        except Exception:
            pass
        return {
            "pdv": nome,
            "cidade": str(r["CIDADE"]).split(" - ")[0][:20],
            "uf": r["UF"],
            "abr": r["abr"],
            "mai_proj": r["mai_proj"],
            "var": r["var_unid"],
        }

    def top_marca(setor, direcao):
        sub = marca_setor_grp[marca_setor_grp["SETOR_NOME"] == setor]
        if len(sub) == 0:
            return None
        if direcao == "alta":
            r = sub.loc[sub["var_unid"].idxmax()]
        else:
            r = sub.loc[sub["var_unid"].idxmin()]
        return {
            "marca": str(r["MARCA"]).replace(" (ALC)", "")[:28],
            "abr": r["abr"],
            "mai_proj": r["mai_proj"],
            "var": r["var_unid"],
        }

    # Top destaques (maior variação positiva)
    destaques_df = setor_grp.sort_values("var_unid", ascending=False).head(top_n)
    destaques = []
    for _, r in destaques_df.iterrows():
        destaques.append({
            "setor": r["SETOR_NOME"],
            "abr": r["abr"],
            "mai_proj": r["mai_proj"],
            "var_unid": r["var_unid"],
            "var_pct": r["var_pct"],
            "pdv_top": top_pdv(r["SETOR_NOME"], "alta"),
            "marca_top": top_marca(r["SETOR_NOME"], "alta"),
        })

    # Top ofensores (menor variação ou queda)
    ofensores_df = setor_grp.sort_values("var_unid", ascending=True).head(top_n)
    ofensores = []
    for _, r in ofensores_df.iterrows():
        ofensores.append({
            "setor": r["SETOR_NOME"],
            "abr": r["abr"],
            "mai_proj": r["mai_proj"],
            "var_unid": r["var_unid"],
            "var_pct": r["var_pct"],
            "pdv_top": top_pdv(r["SETOR_NOME"], "queda"),
            "marca_top": top_marca(r["SETOR_NOME"], "queda"),
        })

    # Agregado por UF do gerente (todas ordenadas por Mai-Proj desc)
    uf_grp = df_gerente.groupby("UF").agg(
        abr=(mes_anterior, "sum"),
        mai_proj=(col_proj, "sum"),
    ).reset_index()
    uf_grp["var_unid"] = uf_grp["mai_proj"] - uf_grp["abr"]
    uf_grp = uf_grp.sort_values("mai_proj", ascending=False).head(8)
    uf_list = [{"uf": r["UF"], "abr": r["abr"], "mai_proj": r["mai_proj"], "var": r["var_unid"]}
               for _, r in uf_grp.iterrows() if pd.notna(r["UF"])]

    # Agregado por Bandeira do gerente (top 8 por Mai-Proj)
    # Junta "OUTRAS" para bandeiras menores
    band_grp = df_gerente.groupby("Bandeira").agg(
        abr=(mes_anterior, "sum"),
        mai_proj=(col_proj, "sum"),
    ).reset_index()
    band_grp["var_unid"] = band_grp["mai_proj"] - band_grp["abr"]
    band_grp = band_grp.sort_values("mai_proj", ascending=False).head(8)
    band_list = [{"band": r["Bandeira"], "abr": r["abr"], "mai_proj": r["mai_proj"], "var": r["var_unid"]}
                 for _, r in band_grp.iterrows()]

    # Série mensal agregada
    serie_mensal = {}
    for m in meses_2026:
        serie_mensal[m] = df_gerente[m].sum()
    serie_mensal[col_proj] = proj_mai

    return {
        "ytd": ytd,
        "proj_mai": proj_mai,
        "abr": abr,
        "var_mai_abr": var_mai_abr,
        "var_mai_abr_pct": var_mai_abr_pct,
        "s1": df_gerente[col_s1].sum(),
        "s2": s2,
        "s3": s3,
        "var_s3_s2": var_s3_s2,
        "var_s3_s2_pct": var_s3_s2_pct,
        "n_setores": df_gerente["SETOR_NOME"].nunique(),
        "n_pdvs": df_gerente["CNPJ"].nunique(),
        "destaques": destaques,
        "ofensores": ofensores,
        "serie_mensal": serie_mensal,
        "uf_list": uf_list,
        "band_list": band_list,
    }


def analise_setor(df_setor, mes_corrente=MES_CORRENTE, mes_anterior=MES_ANTERIOR):
    """Devolve dicionário com todas as métricas e tops do setor."""
    meses_2026 = ["202601", "202602", "202603", "202604"]
    col_proj = f"{mes_corrente}-Proj"
    col_s1 = f"{mes_corrente}-S1"
    col_s2 = f"{mes_corrente}-S2"
    col_s3 = f"{mes_corrente}-S3"

    # Total mensal
    serie_mensal = {}
    for m in meses_2026:
        serie_mensal[m] = df_setor[m].sum()
    serie_mensal[col_proj] = df_setor[col_proj].sum()

    ytd = sum(df_setor[m].sum() for m in meses_2026)
    proj_mai = df_setor[col_proj].sum()
    pdvs_ativos = (df_setor.groupby("CNPJ")[meses_2026 + [col_proj]].sum().sum(axis=1) > 0).sum()

    # Variação Mai-Proj vs Abr
    var_mai_abr = proj_mai - df_setor[mes_anterior].sum()
    var_mai_abr_pct = (var_mai_abr / df_setor[mes_anterior].sum() * 100) if df_setor[mes_anterior].sum() > 0 else 0

    # Semanal: S3 vs S2
    s2 = df_setor[col_s2].sum()
    s3 = df_setor[col_s3].sum()
    var_s3_s2 = s3 - s2
    var_s3_s2_pct = (var_s3_s2 / s2 * 100) if s2 > 0 else 0

    # Por PDV: variação Mai-Proj vs Abr (com threshold)
    pdv_grp = df_setor.groupby(["CNPJ", "PDV", "CIDADE", "UF"]).agg(
        abr=(mes_anterior, "sum"),
        mai_proj=(col_proj, "sum"),
        s2=(col_s2, "sum"),
        s3=(col_s3, "sum"),
    ).reset_index()
    pdv_grp["var_unid"] = pdv_grp["mai_proj"] - pdv_grp["abr"]
    pdv_grp["var_pct"] = np.where(pdv_grp["abr"] > 0,
                                   pdv_grp["var_unid"] / pdv_grp["abr"] * 100, np.nan)

    altas = pdv_grp[pdv_grp["var_unid"] >= THRESHOLD_VARIACAO].sort_values("var_unid", ascending=False).head(5)
    quedas = pdv_grp[pdv_grp["var_unid"] <= -THRESHOLD_VARIACAO].sort_values("var_unid").head(5)

    return {
        "serie_mensal": serie_mensal,
        "ytd": ytd,
        "proj_mai": proj_mai,
        "pdvs_ativos": int(pdvs_ativos),
        "var_mai_abr": var_mai_abr,
        "var_mai_abr_pct": var_mai_abr_pct,
        "s1": df_setor[col_s1].sum(),
        "s2": s2,
        "s3": s3,
        "var_s3_s2": var_s3_s2,
        "var_s3_s2_pct": var_s3_s2_pct,
        "altas": altas,
        "quedas": quedas,
        "total_pdvs_alta": (pdv_grp["var_unid"] >= THRESHOLD_VARIACAO).sum(),
        "total_pdvs_queda": (pdv_grp["var_unid"] <= -THRESHOLD_VARIACAO).sum(),
    }


def analise_produto(df, ytd_cols=YTD_COLS, ytd_ant_cols=YTD_ANT_COLS):
    """Análise consolidada de produto (MARCA + APRESENTACAO)."""
    df = df.copy()
    df["_ytd"] = df[ytd_cols].sum(axis=1)
    df["_ytd_ant"] = df[ytd_ant_cols].sum(axis=1)

    por_marca = df.groupby("MARCA").agg(
        ytd=("_ytd", "sum"),
        ytd_ant=("_ytd_ant", "sum"),
    ).reset_index()
    por_marca["var_unid"] = por_marca["ytd"] - por_marca["ytd_ant"]
    por_marca["var_pct"] = np.where(por_marca["ytd_ant"] > 0,
                                     por_marca["var_unid"] / por_marca["ytd_ant"] * 100, np.nan)
    por_marca = por_marca.sort_values("ytd", ascending=False)

    por_apres = df.groupby("APRESENTACAO").agg(
        ytd=("_ytd", "sum"),
        ytd_ant=("_ytd_ant", "sum"),
    ).reset_index()
    por_apres["var_unid"] = por_apres["ytd"] - por_apres["ytd_ant"]
    por_apres["var_pct"] = np.where(por_apres["ytd_ant"] > 0,
                                     por_apres["var_unid"] / por_apres["ytd_ant"] * 100, np.nan)
    por_apres = por_apres.sort_values("ytd", ascending=False)

    return por_marca, por_apres


# ============================================================
# CHARTS (matplotlib → PNG bytes)
# ============================================================
def chart_evolucao_mensal(serie_mensal, titulo=None):
    """Linha mensal jan-abr + Mai-Proj com destaque na projeção."""
    meses_labels = ["Jan", "Fev", "Mar", "Abr", "Mai-Proj"]
    valores = list(serie_mensal.values())

    fig, ax = plt.subplots(figsize=(7.0, 2.8), dpi=150)
    ax.plot(meses_labels[:4], valores[:4], "-o", color=f"#{NAVY}", linewidth=2.5,
            markersize=8, markerfacecolor=f"#{NAVY}", markeredgecolor="white", markeredgewidth=1.5)
    # Conexão tracejada para projeção
    ax.plot(meses_labels[3:], valores[3:], "--o", color=f"#{GOLD}", linewidth=2.5,
            markersize=8, markerfacecolor=f"#{GOLD}", markeredgecolor="white", markeredgewidth=1.5)

    # Labels nos pontos
    for i, v in enumerate(valores):
        ax.annotate(fmt_num(v), (i, v), textcoords="offset points", xytext=(0, 10),
                    ha="center", fontsize=9, fontweight="bold", color=f"#{NAVY}" if i < 4 else f"#{GOLD}")

    ax.set_ylim(0, max(valores) * 1.25 if max(valores) > 0 else 1)
    ax.grid(axis="y", linestyle=":", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.spines["bottom"].set_color("#cccccc")
    ax.tick_params(colors=f"#{GRAY}", labelsize=10)
    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: fmt_num(x)))
    if titulo:
        ax.set_title(titulo, fontsize=11, color=f"#{NAVY}", fontweight="bold", loc="left", pad=8)

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_semanal(s1, s2, s3):
    """Barras semanais com destaque S3 vs S2."""
    fig, ax = plt.subplots(figsize=(3.5, 2.5), dpi=150)
    labels = ["S1", "S2", "S3"]
    vals = [s1, s2, s3]
    colors = [f"#{LIGHT_GRAY}", f"#{ICE_BLUE}", f"#{NAVY}"]
    bars = ax.bar(labels, vals, color=colors, edgecolor="white", linewidth=1.5, width=0.55)

    for b, v in zip(bars, vals):
        ax.annotate(fmt_num(v), (b.get_x() + b.get_width() / 2, v),
                    textcoords="offset points", xytext=(0, 5), ha="center",
                    fontsize=10, fontweight="bold", color=f"#{NAVY}")

    ax.set_ylim(0, max(vals) * 1.25 if max(vals) > 0 else 1)
    ax.grid(axis="y", linestyle=":", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.spines["bottom"].set_color("#cccccc")
    ax.tick_params(colors=f"#{GRAY}", labelsize=10)
    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: fmt_num(x)))

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_barras_horizontais(df_data, label_col, val_col, color_hex, titulo=None, max_items=8):
    """Barras horizontais para top marcas/apresentações."""
    df_p = df_data.head(max_items).copy().iloc[::-1]  # inverte para top no topo
    fig, ax = plt.subplots(figsize=(5.0, 0.45 * max_items + 0.5), dpi=150)
    bars = ax.barh(df_p[label_col].astype(str), df_p[val_col],
                   color=f"#{color_hex}", edgecolor="white", linewidth=1)

    max_val = df_p[val_col].max() if len(df_p) else 1
    for b, v in zip(bars, df_p[val_col]):
        ax.annotate(fmt_num(v), (v, b.get_y() + b.get_height() / 2),
                    textcoords="offset points", xytext=(5, 0), va="center",
                    fontsize=9, fontweight="bold", color=f"#{NAVY}")

    ax.set_xlim(0, max_val * 1.18 if max_val > 0 else 1)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.tick_params(axis="y", labelsize=9, colors=f"#{NAVY}")
    ax.tick_params(axis="x", labelbottom=False, bottom=False)
    if titulo:
        ax.set_title(titulo, fontsize=11, color=f"#{NAVY}", fontweight="bold", loc="left", pad=8)

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_var_marca(df_marca, max_items=8):
    """Barras horizontais com variação YTD por marca (vermelho/verde).

    Labels sempre fora da barra: positivos à direita da ponta, negativos à
    esquerda da ponta. Os nomes das marcas ficam à esquerda do eixo zero
    (negativo) com margem dedicada, evitando sobreposição.
    """
    df_p = df_marca.head(max_items).copy().iloc[::-1]
    fig, ax = plt.subplots(figsize=(5.0, 0.45 * max_items + 0.5), dpi=150)
    colors = [f"#{ACCENT_GREEN}" if v >= 0 else f"#{ACCENT_RED}" for v in df_p["var_unid"]]
    bars = ax.barh(df_p["MARCA"].astype(str), df_p["var_unid"],
                   color=colors, edgecolor="white", linewidth=1)

    max_abs = max(abs(df_p["var_unid"].min()), abs(df_p["var_unid"].max())) if len(df_p) else 1
    for b, v, p in zip(bars, df_p["var_unid"], df_p["var_pct"]):
        # Labels sempre à DIREITA da ponta (positivo) ou à ESQUERDA da ponta (negativo)
        # nunca dentro da barra para não sobrepor outros elementos
        if v >= 0:
            x_pos = v
            ha = "left"
            offset = 5
        else:
            x_pos = v
            ha = "right"
            offset = -5
        ax.annotate(f"{fmt_num(v)} ({fmt_pct(p)})", (x_pos, b.get_y() + b.get_height() / 2),
                    textcoords="offset points", xytext=(offset, 0), va="center", ha=ha,
                    fontsize=8, fontweight="bold",
                    color=f"#{ACCENT_GREEN}" if v >= 0 else f"#{ACCENT_RED}")

    ax.axvline(0, color=f"#{GRAY}", linewidth=0.8)
    # Folga generosa nos dois lados para acomodar labels longos sem cruzar o eixo
    ax.set_xlim(-max_abs * 2.4 if max_abs > 0 else -1, max_abs * 1.8 if max_abs > 0 else 1)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.tick_params(axis="y", labelsize=9, colors=f"#{NAVY}")
    ax.tick_params(axis="x", labelbottom=False, bottom=False)

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# ============================================================
# PPTX BUILDER
# ============================================================
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # widescreen 16:9


def _set_fill(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(hex_color)
    # Cells não têm atributo .line; só shapes têm.
    if hasattr(shape, "line"):
        try:
            shape.line.fill.background()
        except Exception:
            pass


def _add_text(slide, left, top, width, height, text, *, size=12, bold=False,
              color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return tx


def _kpi_card(slide, left, top, width, height, label, value, *, sub=None,
              accent=NAVY, bg=WHITE, value_color=None):
    """Cartão KPI com label pequeno em cima, valor grande no meio, sub opcional."""
    # Fundo
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.08
    _set_fill(box, bg)
    box.line.color.rgb = hex_to_rgb(LIGHT_GRAY)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False

    # Barra de accent à esquerda
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    _set_fill(bar, accent)

    # Label
    _add_text(slide, left + Inches(0.18), top + Inches(0.08), width - Inches(0.2), Inches(0.25),
              label.upper(), size=8.5, bold=True, color=GRAY)
    # Valor
    vcolor = value_color if value_color else NAVY
    _add_text(slide, left + Inches(0.18), top + Inches(0.32), width - Inches(0.2), Inches(0.55),
              value, size=22, bold=True, color=vcolor)
    # Sub
    if sub:
        sub_color = (ACCENT_GREEN if sub.startswith("+") else
                     ACCENT_RED if sub.startswith("-") and not sub.startswith("—") else GRAY)
        _add_text(slide, left + Inches(0.18), top + height - Inches(0.32),
                  width - Inches(0.2), Inches(0.25),
                  sub, size=9.5, bold=True, color=sub_color)


def add_slide_capa(prs, info):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # Fundo navy
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, NAVY)
    bg.line.fill.background()

    # Accent gold à direita
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(0.18), 0, Inches(0.18), SLIDE_H)
    _set_fill(accent, GOLD)

    _add_text(s, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.5),
              "ALCON | COMMERCIAL INTELLIGENCE", size=14, bold=True, color=ICE_BLUE)

    _add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.5),
              "Análise MDTR por Gerente", size=44, bold=True, color=WHITE, font="Georgia")

    _add_text(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.5),
              f"Período: Abr/25 – Mai/26 (Proj)  •  {info['n_setores']} setores  •  "
              f"{info['n_pdvs']:,} PDVs analisados".replace(",", "."),
              size=16, color=ICE_BLUE)

    _add_text(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(2.2),
              "Performance agregada por gerente  ·  Top setores em alta e ofensores  ·  "
              "Maior PDV impactador em cada setor  ·  Visão consolidada de produto",
              size=14, color=WHITE)

    _add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
              f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}  •  Fonte: MDTRS_FV_TOTAL",
              size=10, color=ICE_BLUE)


def add_slide_setor(prs, setor_nome, gerente, a, n_pdvs_total):
    """Adiciona slide de análise de um setor."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # Header navy
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.75))
    _set_fill(header, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.75), SLIDE_W, Inches(0.05))
    _set_fill(accent, GOLD)

    _add_text(s, Inches(0.4), Inches(0.15), Inches(9.0), Inches(0.45),
              f"Setor: {setor_nome}", size=20, bold=True, color=WHITE, font="Georgia")
    _add_text(s, Inches(0.4), Inches(0.5), Inches(9.0), Inches(0.25),
              f"Gerente: {gerente}", size=10, color=ICE_BLUE)
    _add_text(s, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.4),
              "ANÁLISE MENSAL · MAI vs ABR · S3 vs S2",
              size=9, bold=True, color=ICE_BLUE, align=PP_ALIGN.RIGHT)

    # Linha 1: 4 KPIs
    kpi_top = Inches(1.0)
    kpi_h = Inches(1.1)
    kpi_w = Inches(3.05)
    gap = Inches(0.13)
    left0 = Inches(0.4)

    _kpi_card(s, left0, kpi_top, kpi_w, kpi_h,
              f"YTD 2026 (Jan-Abr) — {TIPO_INFO_ANALISE}",
              fmt_num(a["ytd"]) + " un.",
              sub=f"PDVs ativos: {a['pdvs_ativos']} / {n_pdvs_total}",
              accent=NAVY)

    _kpi_card(s, left0 + (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "MAI-PROJ (Maio Projetado)",
              fmt_num(a["proj_mai"]) + " un.",
              sub=f"{fmt_num(a['var_mai_abr'])} un. vs Abr ({fmt_pct(a['var_mai_abr_pct'])})",
              accent=GOLD,
              value_color=GOLD)

    _kpi_card(s, left0 + 2 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "MAI-S3 vs MAI-S2",
              fmt_num(a["var_s3_s2"]) + " un.",
              sub=fmt_pct(a["var_s3_s2_pct"]),
              accent=ACCENT_GREEN if a["var_s3_s2"] >= 0 else ACCENT_RED,
              value_color=ACCENT_GREEN if a["var_s3_s2"] >= 0 else ACCENT_RED)

    _kpi_card(s, left0 + 3 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "PDVs em DESTAQUE (≥±40)",
              f"↑ {a['total_pdvs_alta']}   ↓ {a['total_pdvs_queda']}",
              sub=f"Threshold: ±{THRESHOLD_VARIACAO} unid.",
              accent=NAVY)

    # Gráfico evolução mensal (esquerda)
    chart_top = Inches(2.3)
    chart_evo = chart_evolucao_mensal(a["serie_mensal"], titulo="Evolução Mensal — Sell-Out (unid.)")
    s.shapes.add_picture(chart_evo, Inches(0.4), chart_top, width=Inches(7.6), height=Inches(2.7))

    # Gráfico semanal (direita)
    _add_text(s, Inches(8.2), chart_top - Inches(0.05), Inches(4.7), Inches(0.3),
              "Desempenho Semanal — Maio/26", size=11, bold=True, color=NAVY)
    chart_sem = chart_semanal(a["s1"], a["s2"], a["s3"])
    s.shapes.add_picture(chart_sem, Inches(8.2), chart_top + Inches(0.15), width=Inches(4.7), height=Inches(2.4))

    # Tops PDVs (abaixo)
    tops_top = Inches(5.15)
    _add_pdv_table(s, Inches(0.4), tops_top, Inches(6.3), a["altas"],
                   "↑ TOP 5 PDVs em ALTA (Mai-Proj vs Abr ≥ +40 unid)", ACCENT_GREEN)
    _add_pdv_table(s, Inches(6.9), tops_top, Inches(6.0), a["quedas"],
                   "↓ TOP 5 PDVs em QUEDA (Mai-Proj vs Abr ≤ −40 unid)", ACCENT_RED)


def _add_pdv_table(slide, left, top, width, df_pdv, titulo, accent):
    """Tabela compacta de PDVs."""
    _add_text(slide, left, top, width, Inches(0.25),
              titulo, size=10.5, bold=True, color=accent)
    table_top = top + Inches(0.3)

    if len(df_pdv) == 0:
        msg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, table_top, width, Inches(1.85))
        _set_fill(msg_box, "F8F8F8")
        msg_box.line.fill.background()
        _add_text(slide, left, table_top + Inches(0.7), width, Inches(0.4),
                  f"Nenhum PDV com variação ≥ ±{THRESHOLD_VARIACAO} unid.",
                  size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return

    # Cabeçalho
    n_rows = min(len(df_pdv), 5) + 1
    row_h = Inches(0.32)
    table = slide.shapes.add_table(n_rows, 4, left, table_top, width, row_h * n_rows).table
    # Larguras
    table.columns[0].width = int(width * 0.55)
    table.columns[1].width = int(width * 0.15)
    table.columns[2].width = int(width * 0.15)
    table.columns[3].width = int(width * 0.15)

    headers = ["PDV — Cidade/UF", "Abr", "Mai-Proj", "Δ unid"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        _set_fill(cell, NAVY)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = h
        r.font.name = "Calibri"
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = hex_to_rgb(WHITE)
        cell.margin_left = Emu(60000)
        cell.margin_right = Emu(60000)
        cell.margin_top = Emu(20000)
        cell.margin_bottom = Emu(20000)

    for i, (_, row) in enumerate(df_pdv.head(5).iterrows(), start=1):
        # Limpa o nome do PDV + sufixo CNPJ para diferenciar lojas da mesma rede
        pdv_name = str(row["PDV"]).split(" - ")[0][:35]
        try:
            suf = str(int(row["CNPJ"]))[-4:]
            pdv_name = f"{pdv_name} ({suf})"
        except Exception:
            pass
        cidade_uf = f"{row['CIDADE']}".replace(f" - {row['UF']}", "")[:25] + f" / {row['UF']}"
        row_color = "F8F8F8" if i % 2 == 0 else WHITE

        cells_data = [
            (f"{pdv_name}  ·  {cidade_uf}", PP_ALIGN.LEFT, NAVY, False),
            (fmt_num(row["abr"]), PP_ALIGN.RIGHT, GRAY, False),
            (fmt_num(row["mai_proj"]), PP_ALIGN.RIGHT, NAVY, True),
            (f"{fmt_num(row['var_unid'])}", PP_ALIGN.RIGHT, accent, True),
        ]
        for j, (txt, align, color, bold) in enumerate(cells_data):
            cell = table.cell(i, j)
            cell.text = ""
            _set_fill(cell, row_color)
            p = cell.text_frame.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.name = "Calibri"
            r.font.size = Pt(9)
            r.font.bold = bold
            r.font.color.rgb = hex_to_rgb(color)
            cell.margin_left = Emu(60000)
            cell.margin_right = Emu(60000)
            cell.margin_top = Emu(15000)
            cell.margin_bottom = Emu(15000)


def add_slide_gerente(prs, gerente, a):
    """Gera 2 slides para o gerente:
       Slide A — Visão Geral: KPIs + evolução mensal + semanal + UF + Bandeira
       Slide B — Top Setores: tabelas com PDV e marca de maior impacto
    """
    _slide_gerente_overview(prs, gerente, a)
    _slide_gerente_setores(prs, gerente, a)


def _slide_gerente_overview(prs, gerente, a):
    """Slide A — Indicadores macro do gerente."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # Header
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.75))
    _set_fill(header, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.75), SLIDE_W, Inches(0.05))
    _set_fill(accent, GOLD)

    _add_text(s, Inches(0.4), Inches(0.15), Inches(9.0), Inches(0.45),
              f"Gerente: {gerente}", size=20, bold=True, color=WHITE, font="Georgia")
    _add_text(s, Inches(0.4), Inches(0.5), Inches(9.0), Inches(0.25),
              f"{a['n_setores']} setores  ·  {a['n_pdvs']:,} PDVs analisados  ·  Visão Geral".replace(",", "."),
              size=10, color=ICE_BLUE)
    _add_text(s, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.4),
              "EVOLUÇÃO · SEMANAL · UF · BANDEIRA",
              size=9, bold=True, color=ICE_BLUE, align=PP_ALIGN.RIGHT)

    # KPIs (linha de 4)
    kpi_top = Inches(1.0)
    kpi_h = Inches(0.95)
    kpi_w = Inches(3.05)
    gap = Inches(0.13)
    left0 = Inches(0.4)

    _kpi_card(s, left0, kpi_top, kpi_w, kpi_h,
              f"YTD 2026 (Jan-Abr) — {TIPO_INFO_ANALISE}",
              fmt_num(a["ytd"]) + " un.",
              sub=f"Abr/26: {fmt_num(a['abr'])} un.",
              accent=NAVY)
    _kpi_card(s, left0 + (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "MAI-PROJ (Maio Projetado)",
              fmt_num(a["proj_mai"]) + " un.",
              sub=f"{fmt_num(a['var_mai_abr'])} un. vs Abr ({fmt_pct(a['var_mai_abr_pct'])})",
              accent=GOLD, value_color=GOLD)
    _kpi_card(s, left0 + 2 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "MAI-S3 vs MAI-S2",
              fmt_num(a["var_s3_s2"]) + " un.",
              sub=fmt_pct(a["var_s3_s2_pct"]),
              accent=ACCENT_GREEN if a["var_s3_s2"] >= 0 else ACCENT_RED,
              value_color=ACCENT_GREEN if a["var_s3_s2"] >= 0 else ACCENT_RED)
    _kpi_card(s, left0 + 3 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "Δ MAI-PROJ vs ABR",
              fmt_pct(a["var_mai_abr_pct"]),
              sub=f"{fmt_num(a['var_mai_abr'])} un.",
              accent=ACCENT_GREEN if a["var_mai_abr"] >= 0 else ACCENT_RED,
              value_color=ACCENT_GREEN if a["var_mai_abr"] >= 0 else ACCENT_RED)

    # Linha 1 dos gráficos: Evolução mensal (esquerda larga) + Semanal (direita)
    g1_top = Inches(2.15)
    g1_h = Inches(2.5)
    chart_evo = chart_evolucao_mensal(a["serie_mensal"],
                                       titulo="Evolução Mensal — Mai/26-Proj")
    s.shapes.add_picture(chart_evo, Inches(0.4), g1_top, width=Inches(8.0), height=g1_h)

    _add_text(s, Inches(8.6), g1_top, Inches(4.3), Inches(0.3),
              "Desempenho Semanal — Maio/26",
              size=11, bold=True, color=NAVY)
    chart_sem = chart_semanal(a["s1"], a["s2"], a["s3"])
    s.shapes.add_picture(chart_sem, Inches(8.6), g1_top + Inches(0.3),
                          width=Inches(4.3), height=Inches(2.2))

    # Linha 2 dos gráficos: UF (esquerda) + Bandeira (direita)
    g2_top = Inches(4.85)
    g2_h = Inches(2.4)

    _add_text(s, Inches(0.4), g2_top, Inches(6.3), Inches(0.3),
              "Desempenho por UF (Mai-Proj)",
              size=11, bold=True, color=NAVY)
    chart_uf = _chart_dim_horiz(a["uf_list"], "uf", titulo=None, max_items=8)
    s.shapes.add_picture(chart_uf, Inches(0.4), g2_top + Inches(0.3),
                          width=Inches(6.3), height=g2_h - Inches(0.3))

    _add_text(s, Inches(6.9), g2_top, Inches(6.0), Inches(0.3),
              "Desempenho por Bandeira (Mai-Proj)",
              size=11, bold=True, color=NAVY)
    chart_band = _chart_dim_horiz(a["band_list"], "band", titulo=None, max_items=8)
    s.shapes.add_picture(chart_band, Inches(6.9), g2_top + Inches(0.3),
                          width=Inches(6.0), height=g2_h - Inches(0.3))


def _slide_gerente_setores(prs, gerente, a):
    """Slide B — Top setores com PDV e MARCA de maior impacto."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # Header
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.75))
    _set_fill(header, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.75), SLIDE_W, Inches(0.05))
    _set_fill(accent, GOLD)

    _add_text(s, Inches(0.4), Inches(0.15), Inches(9.0), Inches(0.45),
              f"Gerente: {gerente}", size=20, bold=True, color=WHITE, font="Georgia")
    _add_text(s, Inches(0.4), Inches(0.5), Inches(9.0), Inches(0.25),
              f"Top {len(a['destaques'])} setores · PDV e Marca de maior impacto",
              size=10, color=ICE_BLUE)
    _add_text(s, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.4),
              "TOP SETORES · PDV · MARCA",
              size=9, bold=True, color=ICE_BLUE, align=PP_ALIGN.RIGHT)

    # 2 tabelas grandes (mais espaço já que é só elas no slide)
    table_top = Inches(1.05)
    _add_setor_pdv_marca_table(s, Inches(0.4), table_top, Inches(6.3), a["destaques"],
                                "↑ TOP SETORES em DESTAQUE — Maior alta (Mai-Proj vs Abr)",
                                ACCENT_GREEN)
    _add_setor_pdv_marca_table(s, Inches(6.9), table_top, Inches(6.0), a["ofensores"],
                                "↓ TOP SETORES OFENSORES — Menor performance (Mai-Proj vs Abr)",
                                ACCENT_RED)


def _chart_dim_horiz(items, key, titulo=None, max_items=8):
    """Barras horizontais p/ UF ou Bandeira: mostra Mai-Proj com variação no rótulo."""
    if not items:
        # Empty placeholder
        fig, ax = plt.subplots(figsize=(5.0, 2.2), dpi=150)
        ax.text(0.5, 0.5, "Sem dados", ha="center", va="center",
                transform=ax.transAxes, color=f"#{GRAY}", fontsize=11, style="italic")
        ax.axis("off")
        buf = io.BytesIO()
        plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf

    items = items[:max_items]
    labels = [str(it[key])[:18] for it in items]
    valores = [it["mai_proj"] for it in items]
    variacoes = [it["var"] for it in items]

    fig, ax = plt.subplots(figsize=(5.0, 0.32 * max_items + 0.3), dpi=150)
    bars = ax.barh(range(len(labels)), valores, color=f"#{NAVY}",
                   edgecolor="white", linewidth=1)
    ax.set_yticks(range(len(labels)))
    ax.set_yticklabels(labels, fontsize=9, color=f"#{NAVY}")
    ax.invert_yaxis()  # top first

    max_val = max(valores) if valores else 1
    for i, (v, var) in enumerate(zip(valores, variacoes)):
        var_color = f"#{ACCENT_GREEN}" if var >= 0 else f"#{ACCENT_RED}"
        sign = "+" if var >= 0 else ""
        label_text = f"{fmt_num(v)}  ({sign}{fmt_num(var)})"
        ax.text(v + max_val * 0.015, i, label_text,
                va="center", fontsize=8, fontweight="bold",
                color=f"#{NAVY}")
        # Pequeno marker colorido após o valor
        # (alternativa: usar 2 labels — mantemos só um e cor neutra para legibilidade)

    ax.set_xlim(0, max_val * 1.35)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.tick_params(axis="x", labelbottom=False, bottom=False)

    if titulo:
        ax.set_title(titulo, fontsize=11, color=f"#{NAVY}", fontweight="bold", loc="left", pad=8)

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_setor_pdv_marca_table(slide, left, top, width, items, titulo, accent):
    """Tabela 3-linhas-por-item: setor + ↳ PDV de maior impacto + ↳ Marca de maior impacto."""
    _add_text(slide, left, top, width, Inches(0.28),
              titulo, size=11, bold=True, color=accent)
    table_top = top + Inches(0.32)

    if not items or len(items) == 0:
        msg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, table_top, width, Inches(5.0))
        _set_fill(msg_box, "F8F8F8")
        msg_box.line.fill.background()
        _add_text(slide, left, table_top + Inches(2.0), width, Inches(0.4),
                  "Sem dados para este gerente.",
                  size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return

    n_items = min(len(items), 5)
    n_rows = 1 + n_items * 3   # header + 3 linhas por item
    row_h = Inches(0.27)
    table = slide.shapes.add_table(n_rows, 4, left, table_top, width, row_h * n_rows).table

    table.columns[0].width = int(width * 0.55)
    table.columns[1].width = int(width * 0.15)
    table.columns[2].width = int(width * 0.15)
    table.columns[3].width = int(width * 0.15)

    headers = ["Setor  ·  ↳ PDV  ·  ↳ Marca", "Abr", "Mai-Proj", "Δ unid"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        _set_fill(cell, NAVY)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = h
        r.font.name = "Calibri"; r.font.size = Pt(9); r.font.bold = True
        r.font.color.rgb = hex_to_rgb(WHITE)
        cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
        cell.margin_top = Emu(15000); cell.margin_bottom = Emu(15000)

    for i, item in enumerate(items[:5]):
        row_idx_setor = 1 + i * 3
        row_idx_pdv = row_idx_setor + 1
        row_idx_marca = row_idx_setor + 2
        bg_main = "F0F2F5" if i % 2 == 0 else WHITE
        bg_sub = "FAFAFA" if i % 2 == 0 else "F8F8F8"

        # LINHA 1: SETOR (negrito, navy)
        setor_var_color = ACCENT_GREEN if item["var_unid"] >= 0 else ACCENT_RED
        _fill_row(table, row_idx_setor, [
            (item["setor"][:42], PP_ALIGN.LEFT, NAVY, True, 9.5, False),
            (fmt_num(item["abr"]), PP_ALIGN.RIGHT, GRAY, False, 9, False),
            (fmt_num(item["mai_proj"]), PP_ALIGN.RIGHT, NAVY, True, 9, False),
            (fmt_num(item["var_unid"]), PP_ALIGN.RIGHT, setor_var_color, True, 9, False),
        ], bg_main, top_pad=12000)

        # LINHA 2: PDV
        pdv = item["pdv_top"]
        if pdv is not None:
            pdv_var_color = ACCENT_GREEN if pdv["var"] >= 0 else ACCENT_RED
            pdv_text = f"   ↳ PDV: {pdv['pdv']}  ·  {pdv['cidade']}/{pdv['uf']}"
            _fill_row(table, row_idx_pdv, [
                (pdv_text[:60], PP_ALIGN.LEFT, GRAY, False, 8.5, True),
                (fmt_num(pdv["abr"]), PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                (fmt_num(pdv["mai_proj"]), PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                (fmt_num(pdv["var"]), PP_ALIGN.RIGHT, pdv_var_color, True, 8.5, True),
            ], bg_sub, top_pad=6000)
        else:
            _fill_row(table, row_idx_pdv,
                      [("   ↳ PDV: —", PP_ALIGN.LEFT, GRAY, False, 8.5, True),
                       ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                       ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                       ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5, True)],
                      bg_sub, top_pad=6000)

        # LINHA 3: MARCA
        marca = item["marca_top"]
        if marca is not None:
            marca_var_color = ACCENT_GREEN if marca["var"] >= 0 else ACCENT_RED
            marca_text = f"   ↳ Produto: {marca['marca']}"
            _fill_row(table, row_idx_marca, [
                (marca_text[:55], PP_ALIGN.LEFT, GRAY, False, 8.5, True),
                (fmt_num(marca["abr"]), PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                (fmt_num(marca["mai_proj"]), PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                (fmt_num(marca["var"]), PP_ALIGN.RIGHT, marca_var_color, True, 8.5, True),
            ], bg_sub, top_pad=6000)
        else:
            _fill_row(table, row_idx_marca,
                      [("   ↳ Produto: —", PP_ALIGN.LEFT, GRAY, False, 8.5, True),
                       ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                       ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5, True),
                       ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5, True)],
                      bg_sub, top_pad=6000)


def _fill_row(table, row_idx, cells_data, bg_color, top_pad=12000):
    """Helper para preencher uma linha da tabela com cells_data:
       lista de tuplas (text, align, color_hex, bold, font_size, italic)."""
    for j, (txt, align, color, bold, size, italic) in enumerate(cells_data):
        cell = table.cell(row_idx, j)
        cell.text = ""
        _set_fill(cell, bg_color)
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = hex_to_rgb(color)
        cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
        cell.margin_top = Emu(top_pad); cell.margin_bottom = Emu(top_pad)


def _add_setor_pdv_table(slide, left, top, width, items, titulo, accent):
    """Tabela duas-linhas-por-item: setor (com totais) + PDV maior impacto (recuado, em cinza)."""
    _add_text(slide, left, top, width, Inches(0.25),
              titulo, size=10.5, bold=True, color=accent)
    table_top = top + Inches(0.3)

    if not items or len(items) == 0:
        msg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, table_top, width, Inches(3.0))
        _set_fill(msg_box, "F8F8F8")
        msg_box.line.fill.background()
        _add_text(slide, left, table_top + Inches(1.2), width, Inches(0.4),
                  "Sem dados para este gerente.",
                  size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return

    # Cada item ocupa 2 linhas (setor + PDV). +1 header
    n_items = min(len(items), 5)
    n_rows = 1 + n_items * 2
    row_h = Inches(0.28)
    table = slide.shapes.add_table(n_rows, 4, left, table_top, width, row_h * n_rows).table

    table.columns[0].width = int(width * 0.55)
    table.columns[1].width = int(width * 0.15)
    table.columns[2].width = int(width * 0.15)
    table.columns[3].width = int(width * 0.15)

    # Header
    headers = ["Setor  ·  PDV de maior impacto", "Abr", "Mai-Proj", "Δ unid"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        _set_fill(cell, NAVY)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = h
        r.font.name = "Calibri"
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = hex_to_rgb(WHITE)
        cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
        cell.margin_top = Emu(15000); cell.margin_bottom = Emu(15000)

    for i, item in enumerate(items[:5]):
        row_idx_setor = 1 + i * 2
        row_idx_pdv = row_idx_setor + 1
        # Cor alternada por par (setor+pdv)
        bg_setor = "F0F2F5" if i % 2 == 0 else WHITE
        bg_pdv = "FAFAFA" if i % 2 == 0 else "F8F8F8"

        # LINHA 1: SETOR
        setor_var_color = ACCENT_GREEN if item["var_unid"] >= 0 else ACCENT_RED
        cells_setor = [
            (item["setor"][:42], PP_ALIGN.LEFT, NAVY, True, 9.5),
            (fmt_num(item["abr"]), PP_ALIGN.RIGHT, GRAY, False, 9),
            (fmt_num(item["mai_proj"]), PP_ALIGN.RIGHT, NAVY, True, 9),
            (fmt_num(item["var_unid"]), PP_ALIGN.RIGHT, setor_var_color, True, 9),
        ]
        for j, (txt, align, color, bold, size) in enumerate(cells_setor):
            cell = table.cell(row_idx_setor, j)
            cell.text = ""
            _set_fill(cell, bg_setor)
            p = cell.text_frame.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = hex_to_rgb(color)
            cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
            cell.margin_top = Emu(12000); cell.margin_bottom = Emu(12000)

        # LINHA 2: PDV de maior impacto
        pdv = item["pdv_top"]
        if pdv is not None:
            pdv_var_color = ACCENT_GREEN if pdv["var"] >= 0 else ACCENT_RED
            pdv_text = f"   ↳ {pdv['pdv']}  ·  {pdv['cidade']}/{pdv['uf']}"
            cells_pdv = [
                (pdv_text[:60], PP_ALIGN.LEFT, GRAY, False, 8.5),
                (fmt_num(pdv["abr"]), PP_ALIGN.RIGHT, GRAY, False, 8.5),
                (fmt_num(pdv["mai_proj"]), PP_ALIGN.RIGHT, GRAY, False, 8.5),
                (fmt_num(pdv["var"]), PP_ALIGN.RIGHT, pdv_var_color, True, 8.5),
            ]
        else:
            cells_pdv = [
                ("   ↳ sem PDV identificado", PP_ALIGN.LEFT, GRAY, False, 8.5),
                ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5),
                ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5),
                ("—", PP_ALIGN.RIGHT, GRAY, False, 8.5),
            ]

        for j, (txt, align, color, bold, size) in enumerate(cells_pdv):
            cell = table.cell(row_idx_pdv, j)
            cell.text = ""
            _set_fill(cell, bg_pdv)
            p = cell.text_frame.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = True
            r.font.color.rgb = hex_to_rgb(color)
            cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
            cell.margin_top = Emu(8000); cell.margin_bottom = Emu(8000)


def add_slide_produto(prs, por_marca, por_apres, ytd_label, ytd_ant_label):
    """Slide consolidado de produto: MARCA + APRESENTACAO."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # Header
    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.75))
    _set_fill(header, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.75), SLIDE_W, Inches(0.05))
    _set_fill(accent, GOLD)
    _add_text(s, Inches(0.4), Inches(0.15), Inches(9.0), Inches(0.45),
              "Análise de Produto · Visão Consolidada",
              size=20, bold=True, color=WHITE, font="Georgia")
    _add_text(s, Inches(0.4), Inches(0.5), Inches(9.0), Inches(0.25),
              f"{ytd_label} vs {ytd_ant_label}", size=10, color=ICE_BLUE)
    _add_text(s, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.4),
              "MARCA · APRESENTAÇÃO · YTD",
              size=9, bold=True, color=ICE_BLUE, align=PP_ALIGN.RIGHT)

    # KPIs totais
    total_ytd = por_marca["ytd"].sum()
    total_ant = por_marca["ytd_ant"].sum()
    var_total = total_ytd - total_ant
    var_total_pct = (var_total / total_ant * 100) if total_ant > 0 else 0
    melhor_marca = por_marca.loc[por_marca["var_unid"].idxmax(), "MARCA"] if len(por_marca) else "—"
    pior_marca = por_marca.loc[por_marca["var_unid"].idxmin(), "MARCA"] if len(por_marca) else "—"
    n_marcas_alta = (por_marca["var_unid"] > 0).sum()

    kpi_top = Inches(1.0)
    kpi_h = Inches(1.1)
    kpi_w = Inches(3.05)
    gap = Inches(0.13)
    left0 = Inches(0.4)

    _kpi_card(s, left0, kpi_top, kpi_w, kpi_h,
              ytd_label, fmt_num(total_ytd) + " un.",
              sub=f"vs {ytd_ant_label}: {fmt_num(total_ant)} un.",
              accent=NAVY)
    _kpi_card(s, left0 + (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "Variação Total",
              fmt_num(var_total) + " un.",
              sub=fmt_pct(var_total_pct),
              accent=ACCENT_GREEN if var_total >= 0 else ACCENT_RED,
              value_color=ACCENT_GREEN if var_total >= 0 else ACCENT_RED)
    _kpi_card(s, left0 + 2 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "Marca em MAIOR ALTA",
              str(melhor_marca).replace(" (ALC)", "")[:22],
              sub=f"{fmt_num(por_marca['var_unid'].max())} un.",
              accent=ACCENT_GREEN,
              value_color=ACCENT_GREEN)
    _kpi_card(s, left0 + 3 * (kpi_w + gap), kpi_top, kpi_w, kpi_h,
              "Marca em MAIOR QUEDA",
              str(pior_marca).replace(" (ALC)", "")[:22],
              sub=f"{fmt_num(por_marca['var_unid'].min())} un.",
              accent=ACCENT_RED,
              value_color=ACCENT_RED)

    # Gráficos: marca à esquerda, apresentação à direita
    g_top = Inches(2.3)

    # Marca: ranking YTD
    _add_text(s, Inches(0.4), g_top - Inches(0.05), Inches(6.2), Inches(0.3),
              f"Ranking por MARCA — {ytd_label}",
              size=11, bold=True, color=NAVY)
    por_marca_clean = por_marca.copy()
    por_marca_clean["MARCA"] = por_marca_clean["MARCA"].str.replace(" (ALC)", "", regex=False)
    chart_m = chart_barras_horizontais(por_marca_clean, "MARCA", "ytd", NAVY, max_items=8)
    s.shapes.add_picture(chart_m, Inches(0.4), g_top + Inches(0.2), width=Inches(6.2), height=Inches(4.6))

    # Apresentação: ranking YTD
    _add_text(s, Inches(6.8), g_top - Inches(0.05), Inches(6.2), Inches(0.3),
              f"Variação por MARCA — {ytd_label} vs {ytd_ant_label}",
              size=11, bold=True, color=NAVY)
    chart_var = chart_var_marca(por_marca_clean, max_items=8)
    s.shapes.add_picture(chart_var, Inches(6.8), g_top + Inches(0.2), width=Inches(6.2), height=Inches(4.6))


def add_slide_apresentacao(prs, por_apres, ytd_label, ytd_ant_label):
    """Slide adicional: análise por apresentação."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.75))
    _set_fill(header, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.75), SLIDE_W, Inches(0.05))
    _set_fill(accent, GOLD)
    _add_text(s, Inches(0.4), Inches(0.15), Inches(9.0), Inches(0.45),
              "Análise de Produto · Por Apresentação",
              size=20, bold=True, color=WHITE, font="Georgia")
    _add_text(s, Inches(0.4), Inches(0.5), Inches(9.0), Inches(0.25),
              f"{ytd_label} vs {ytd_ant_label}", size=10, color=ICE_BLUE)

    # Tabela de apresentação
    n_show = min(len(por_apres), 15)
    table_top = Inches(1.0)
    width = Inches(12.5)
    row_h = Inches(0.32)
    table = s.shapes.add_table(n_show + 1, 5, Inches(0.4), table_top, width, row_h * (n_show + 1)).table
    table.columns[0].width = int(width * 0.46)
    table.columns[1].width = int(width * 0.135)
    table.columns[2].width = int(width * 0.135)
    table.columns[3].width = int(width * 0.135)
    table.columns[4].width = int(width * 0.135)

    headers = ["APRESENTAÇÃO", ytd_ant_label, ytd_label, "Δ unid", "Δ %"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        _set_fill(cell, NAVY)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = h
        r.font.name = "Calibri"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = hex_to_rgb(WHITE)
        cell.margin_left = Emu(80000); cell.margin_right = Emu(80000)
        cell.margin_top = Emu(25000); cell.margin_bottom = Emu(25000)

    for i, (_, row) in enumerate(por_apres.head(n_show).iterrows(), start=1):
        row_color = "F8F8F8" if i % 2 == 0 else WHITE
        var_color = ACCENT_GREEN if row["var_unid"] >= 0 else ACCENT_RED
        cells_data = [
            (str(row["APRESENTACAO"])[:55], PP_ALIGN.LEFT, NAVY, False),
            (fmt_num(row["ytd_ant"]), PP_ALIGN.RIGHT, GRAY, False),
            (fmt_num(row["ytd"]), PP_ALIGN.RIGHT, NAVY, True),
            (fmt_num(row["var_unid"]), PP_ALIGN.RIGHT, var_color, True),
            (fmt_pct(row["var_pct"]), PP_ALIGN.RIGHT, var_color, True),
        ]
        for j, (txt, align, color, bold) in enumerate(cells_data):
            cell = table.cell(i, j)
            cell.text = ""
            _set_fill(cell, row_color)
            p = cell.text_frame.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.name = "Calibri"
            r.font.size = Pt(9.5)
            r.font.bold = bold
            r.font.color.rgb = hex_to_rgb(color)
            cell.margin_left = Emu(80000); cell.margin_right = Emu(80000)
            cell.margin_top = Emu(18000); cell.margin_bottom = Emu(18000)


# ============================================================
# MAIN
# ============================================================
def main(input_path=None, output_path=None):
    input_path = input_path or DEFAULT_INPUT
    output_path = output_path or DEFAULT_OUTPUT
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    df = load_data(input_path)

    print(f"[2/5] Filtrando TIPO_INFORMACAO = '{TIPO_INFO_ANALISE}' ...")
    df_filtro = df[df["TIPO_INFORMACAO"] == TIPO_INFO_ANALISE].copy()
    if EXCLUIR_NAO_VISITADO:
        df_filtro = df_filtro[df_filtro["SETOR_NOME"] != "NÃO VISITADO"]
    print(f"      Linhas filtradas: {len(df_filtro):,}")

    print(f"[3/5] Listando gerentes ...")
    period_cols = YTD_COLS + [f"{MES_CORRENTE}-Proj"]
    df_filtro["_total_anal"] = df_filtro[period_cols].sum(axis=1)
    gerentes_validos = (df_filtro.groupby("GERENTE")["_total_anal"].sum()
                        .reset_index()
                        .query("_total_anal > 0")
                        .sort_values("GERENTE"))
    print(f"      Gerentes a gerar: {len(gerentes_validos)}")

    print(f"[4/5] Construindo apresentação ...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Capa
    add_slide_capa(prs, {
        "n_setores": df_filtro["SETOR_NOME"].nunique(),
        "n_pdvs": int(df_filtro["CNPJ"].nunique()),
    })

    # Slides produto
    por_marca, por_apres = analise_produto(df_filtro)
    add_slide_produto(prs, por_marca, por_apres, YTD_LABEL, YTD_ANT_LABEL)
    add_slide_apresentacao(prs, por_apres, YTD_LABEL, YTD_ANT_LABEL)

    # 1 slide por gerente
    for _, grow in gerentes_validos.iterrows():
        gerente = grow["GERENTE"]
        df_g = df_filtro[df_filtro["GERENTE"] == gerente]
        a = analise_gerente(df_g)
        add_slide_gerente(prs, gerente, a)
        print(f"      ✓ {gerente[:30]:30s}  YTD={fmt_num(a['ytd'])}  setores={a['n_setores']}  PDVs={a['n_pdvs']}")

    print(f"[5/5] Salvando em {output_path} ...")
    prs.save(output_path)
    print(f"      ✓ Pronto! {len(prs.slides)} slides gerados.")
    return output_path


if __name__ == "__main__":
    inp = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_INPUT
    out = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_OUTPUT
    main(inp, out)
